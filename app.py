import streamlit as st
import numpy as np
from dotenv import load_dotenv
import os
import pickle
from PyPDF2 import PdfReader
from sentence_transformers import SentenceTransformer
import faiss
from transformers import pipeline
from langchain.text_splitter import RecursiveCharacterTextSplitter
from streamlit_extras import *
from pptx import Presentation
from docx import Document
from concurrent.futures import ThreadPoolExecutor

# Initialize the SentenceTransformer model for embeddings
embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
summarization_model = pipeline("summarization", model="t5-small", tokenizer="t5-small")
qa_model = pipeline("question-answering", model="distilbert-base-uncased-distilled-squad")

def extract_text_from_pdf(pdf):
    pdf_reader = PdfReader(pdf)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def extract_text_from_ppt(ppt):
    presentation = Presentation(ppt)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text

def extract_text_from_docx(docx):
    document = Document(docx)
    text = ""
    for para in document.paragraphs:
        text += para.text
    return text

def summarize_text(text):
    max_length = 512
    if len(text) > max_length:
        text = text[:max_length]
    summary = summarization_model(text, max_length=300, min_length=30, do_sample=False)
    return summary[0]['summary_text']

def create_embeddings(chunks):
    with ThreadPoolExecutor() as executor:
        chunk_embeddings = list(executor.map(lambda chunk: embedding_model.encode(chunk, convert_to_numpy=True), chunks))
    return chunk_embeddings

def main():
    st.header("Gen AI Digital Assistant")
    st.markdown('<p style="font-size:20px">By: Aditi, Mrudula, Risha</p>', unsafe_allow_html=True)

    uploaded_file = st.file_uploader("Upload your PDF, PPT, or Word document", type=['pdf', 'pptx', 'docx'])

    if uploaded_file is not None:
        file_extension = uploaded_file.name.split('.')[-1]

        if file_extension == 'pdf':
            text = extract_text_from_pdf(uploaded_file)
        elif file_extension == 'pptx':
            text = extract_text_from_ppt(uploaded_file)
        elif file_extension == 'docx':
            text = extract_text_from_docx(uploaded_file)
        else:
            st.error("Unsupported file type.")
            return

        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
        chunks = text_splitter.split_text(text=text)

        store_name = uploaded_file.name[:-4]
        st.write(f'{store_name}')

        if os.path.exists(f"{store_name}.pkl"):
            with open(f"{store_name}.pkl", "rb") as f:
                index, texts = pickle.load(f)
        else:
            chunk_embeddings = create_embeddings(chunks)
            chunk_embeddings = np.array(chunk_embeddings)

            index = faiss.IndexFlatL2(chunk_embeddings.shape[1])
            index.add(chunk_embeddings)

            with open(f"{store_name}.pkl", "wb") as f:
                pickle.dump((index, chunks), f)

        if st.button("Summarize Document"):
            summary = summarize_text(text)
            st.write("### Document Summary")
            st.write(summary)

        query = st.text_input("Ask questions about your document:")

        if query:
            query_embedding = embedding_model.encode([query], convert_to_numpy=True)
            distances, indices = index.search(query_embedding, k=3)
            docs = [chunks[i] for i in indices[0]]

            responses = [qa_model(question=query, context=doc) for doc in docs]

            for response in responses:
                st.write(response['answer'])

if __name__ == '__main__':
    main()
